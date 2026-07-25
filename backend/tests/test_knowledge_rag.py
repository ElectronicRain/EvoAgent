from __future__ import annotations

import io
import json
import sqlite3
import uuid

import pytest
from docx import Document
from pptx import Presentation

from backend.app.models import KnowledgeProviderConfig
from backend.app.services.knowledge_processing import (
    ExtractedSection,
    clean_text,
    extract_sections,
    hierarchical_chunks,
)
from backend.app.services.knowledge_vector import EmbeddingClient, RerankClient
from backend.app.services.secrets import secret_store


def _new_base(client):
    response = client.post(
        "/api/knowledge-bases",
        json={"name": f"测试知识库-{uuid.uuid4()}", "discipline": "计算机", "description": "RAG 测试"},
    )
    assert response.status_code == 201
    return response.json()


def test_cleaning_and_hierarchical_chunking():
    cleaned, stats = clean_text("页眉\n页眉\n页眉\n\n正文 A。\n正文 B。\n\n页眉")
    assert cleaned.count("页眉") == 1
    assert stats["repeated_lines_removed"] == 3
    drafts, chunk_stats = hierarchical_chunks(
        [ExtractedSection(("网格质量评价包括正交性、扭曲度和长宽比。 " * 80), "评价指标", "第 2 页")]
    )
    assert any(item.level == "parent" for item in drafts)
    assert any(item.level == "child" and item.parent_index is not None for item in drafts)
    assert all(item.metadata["locator"] == "第 2 页" for item in drafts)
    assert chunk_stats["duplicate_chunks_removed"] >= 0


def test_numbered_list_chunking_preserves_complete_structure():
    content = (
        "学科知识平台建设包含以下五点：\n"
        "1. 建立统一的数据标准。\n"
        "2. 汇聚经过审核的学科资料。\n"
        "3. 建设可追溯的混合检索链路。\n"
        "4. 引入专家复核与安全治理。\n"
        "5. 持续开展质量评估与迭代。"
    )
    drafts, stats = hierarchical_chunks([ExtractedSection(content, "建设方案", "第 3 页")])
    list_parents = [
        item
        for item in drafts
        if item.level == "parent" and item.metadata.get("structure") == "numbered_list"
    ]
    list_children = [
        item
        for item in drafts
        if item.level == "child" and item.metadata.get("structure") == "numbered_list"
    ]
    assert len(list_parents) == 1
    assert "1. 建立统一" in list_parents[0].content
    assert "5. 持续开展" in list_parents[0].content
    assert {item.metadata["list_item_index"] for item in list_children} == {1, 2, 3, 4, 5}
    assert all(item.metadata["list_item_count"] == 5 for item in list_children)
    assert stats["numbered_lists"] == 1
    assert stats["numbered_list_items"] == 5


def test_numbered_list_reset_creates_a_new_focused_list():
    content = (
        "前三章复习要点：\n"
        "1. 进程管理。\n"
        "2. 处理机调度。\n"
        "3. 死锁处理。\n"
        "4. 内存管理。\n"
        "虚拟内存用于扩展可用地址空间，具有以下五点：\n"
        "1. 解决物理内存容量不足。\n"
        "2. 实现进程地址空间隔离。\n"
        "3. 简化内存管理并支持共享。\n"
        "4. 程序编译时使用连续虚拟地址。\n"
        "5. 按需加载页面。"
    )
    drafts, stats = hierarchical_chunks([ExtractedSection(content, "操作系统", "第 24 页")])
    list_parents = [
        item
        for item in drafts
        if item.level == "parent" and item.metadata.get("structure") == "numbered_list"
    ]
    assert [item.metadata["list_item_count"] for item in list_parents] == [4, 5]
    assert "虚拟内存" in list_parents[1].content
    assert "1. 解决物理内存容量不足" in list_parents[1].content
    assert "5. 按需加载页面" in list_parents[1].content
    assert stats["numbered_lists"] == 2
    assert stats["numbered_list_items"] == 9


def test_cross_page_numbered_list_is_stitched_before_chunking():
    sections = [
        ExtractedSection(
            (
                "虚拟内存具有以下五点：\n"
                "1. 解决物理内存容量不足。\n"
                "2. 实现进程地址空间隔离。\n"
                "3. 简化内存管理并支持共享。"
            ),
            "内存管理",
            "第 24 页",
            {"page": 24},
        ),
        ExtractedSection(
            (
                "4. 程序编译时使用连续虚拟地址。\n"
                "5. 按需加载页面。\n"
                "实现原理包括：\n"
                "1. 时间局部性。\n"
                "2. 空间局部性。"
            ),
            "内存管理",
            "第 25 页",
            {"page": 25},
        ),
    ]
    drafts, stats = hierarchical_chunks(sections)
    list_parents = [
        item
        for item in drafts
        if item.level == "parent" and item.metadata.get("structure") == "numbered_list"
    ]
    virtual_memory = next(item for item in list_parents if "虚拟内存" in item.content)
    assert virtual_memory.metadata["list_item_count"] == 5
    assert virtual_memory.metadata["cross_section_continuation"] is True
    assert virtual_memory.metadata["page_start"] == 24
    assert virtual_memory.metadata["page_end"] == 25
    assert virtual_memory.metadata["locators"] == ["第 24 页", "第 25 页"]
    assert "1. 解决物理内存容量不足" in virtual_memory.content
    assert "5. 按需加载页面" in virtual_memory.content
    assert any(item.metadata["list_item_count"] == 2 for item in list_parents)
    assert stats["numbered_lists"] == 2
    assert stats["numbered_list_items"] == 7


def test_pptx_extraction():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "结构化网格"
    slide.placeholders[1].text = "正交性与雅可比行列式"
    buffer = io.BytesIO()
    presentation.save(buffer)
    sections, mime = extract_sections("lecture.pptx", buffer.getvalue())
    assert mime.endswith("presentationml.presentation")
    assert sections[0].metadata["slide"] == 1
    assert "雅可比" in sections[0].text


def test_docx_automatic_numbering_is_preserved_for_list_chunking():
    document = Document()
    document.add_heading("建设方案", level=1)
    document.add_paragraph("平台建设包含以下五点：")
    for value in ("统一标准", "汇聚资料", "混合检索", "专家复核", "持续评估"):
        document.add_paragraph(value, style="List Number")
    buffer = io.BytesIO()
    document.save(buffer)

    sections, _mime = extract_sections("plan.docx", buffer.getvalue())
    assert "1. 统一标准" in sections[0].text
    assert "5. 持续评估" in sections[0].text
    drafts, stats = hierarchical_chunks(sections)
    assert stats["numbered_lists"] == 1
    assert stats["numbered_list_items"] == 5
    assert any(
        item.level == "parent"
        and item.metadata.get("structure") == "numbered_list"
        and "5. 持续评估" in item.content
        for item in drafts
    )


def test_provider_config_is_safe_and_has_defaults(client):
    response = client.get("/api/knowledge/config")
    assert response.status_code == 200
    data = response.json()
    assert data["embedding_model"] == "Qwen/Qwen3-VL-Embedding-8B"
    assert data["rerank_model"] == "BAAI/bge-reranker-v2-m3"
    assert "api_key_ciphertext" not in data


def test_full_text_ingestion_and_rag_query(client):
    base = _new_base(client)
    content = (
        "二维结构化网格质量通常通过正交性、长宽比、扭曲度和雅可比行列式评价。\n\n"
        "雅可比行列式为负意味着单元发生翻转，应当判定为无效网格。\n\n"
        "评价时还应结合目标数值求解器进行误差和稳定性验证。"
    )
    uploaded = client.post(
        f"/api/knowledge-bases/{base['id']}/documents/text",
        json={"title": "网格质量指南", "content": content, "source": "测试标准"},
    )
    assert uploaded.status_code == 201
    document_id = uploaded.json()["id"]

    overview = client.get(f"/api/knowledge-bases/{base['id']}/overview")
    assert overview.status_code == 200
    overview_data = overview.json()
    assert overview_data["statistics"]["parent_chunks"] >= 1
    assert overview_data["statistics"]["child_chunks"] >= 1
    assert overview_data["statistics"]["embeddings"] >= 1
    assert overview_data["retrieval_strategy"]["fusion"].startswith("Reciprocal Rank Fusion")

    detail = client.get(f"/api/knowledge-documents/{document_id}")
    assert detail.status_code == 200
    assert "雅可比" in detail.json()["cleaned_content"]
    assert detail.json()["cleaning_stats"]["cleaned_chars"] > 0

    chunks = client.get(f"/api/knowledge-documents/{document_id}/chunks?level=child")
    assert chunks.status_code == 200
    assert chunks.json()["total"] >= 1
    assert all(item["level"] == "child" for item in chunks.json()["items"])
    assert all(item["embedding"]["indexed"] for item in chunks.json()["items"])

    reindexed = client.post(f"/api/knowledge-bases/{base['id']}/reindex")
    assert reindexed.status_code == 200
    assert reindexed.json()["embedded_chunks"] == chunks.json()["total"]
    refreshed = client.get(f"/api/knowledge-bases/{base['id']}/overview").json()
    assert refreshed["statistics"]["embeddings"] == chunks.json()["total"]

    response = client.post(
        "/api/knowledge/query",
        json={
            "query": "怎样判断二维结构化网格是否发生单元翻转？",
            "knowledge_base_ids": [base["id"]],
            "top_k": 4,
            "generate_answer": True,
        },
    )
    assert response.status_code == 200, response.text
    result = response.json()
    assert result["chunks"]
    assert any("雅可比" in item["content"] for item in result["chunks"])
    assert result["citations"][0]["document_id"]
    assert result["trace"]["embedding_model"] == "Qwen/Qwen3-VL-Embedding-8B"
    assert result["answer"]


def test_exhaustive_query_expands_complete_numbered_list(client):
    base = _new_base(client)
    content = (
        "学科知识平台建设包含以下五点：\n"
        "1. 建立统一的数据标准和目录规范。\n"
        "2. 汇聚经过审核的教材、论文和行业标准。\n"
        "3. 建设支持引用追踪的混合检索链路。\n"
        "4. 引入专家复核、权限控制和安全治理。\n"
        "5. 持续开展质量评估、反馈收集和版本迭代。"
    )
    uploaded = client.post(
        f"/api/knowledge-bases/{base['id']}/documents/text",
        json={"title": "五点建设方案", "content": content, "source": "完整性测试"},
    )
    assert uploaded.status_code == 201
    document_id = uploaded.json()["id"]
    chunk_rows = client.get(
        f"/api/knowledge-documents/{document_id}/chunks?level=all&limit=100"
    ).json()["items"]
    list_parent = next(
        item
        for item in chunk_rows
        if item["level"] == "parent" and item["metadata"].get("structure") == "numbered_list"
    )
    assert list_parent["metadata"]["list_item_count"] == 5
    assert "5. 持续开展" in list_parent["content"]

    response = client.post(
        "/api/knowledge/query",
        json={
            "query": "学科知识平台建设包含哪五点？请完整列出。",
            "knowledge_base_ids": [base["id"]],
            "top_k": 6,
            "generate_answer": False,
        },
    )
    assert response.status_code == 200, response.text
    result = response.json()
    assert result["chunks"]
    assert any("1. 建立统一" in item["context"] for item in result["chunks"])
    assert any("5. 持续开展" in item["context"] for item in result["chunks"])
    assert result["trace"]["exhaustive_query"] is True
    assert result["trace"]["per_document_limit"] == 6
    assert result["trace"]["list_contexts"][0]["item_count"] == 5


def test_streamed_rag_query_reports_explicit_steps_in_order(client):
    base = _new_base(client)
    created = client.post(
        f"/api/knowledge-bases/{base['id']}/documents/text",
        json={
            "title": "流式检索资料",
            "content": "混合检索先进行查询改写，再执行向量与关键词召回、融合和重排序。" * 12,
            "source": "流式测试",
        },
    )
    assert created.status_code == 201

    response = client.post(
        "/api/knowledge/query/stream",
        json={
            "query": "混合检索如何工作？",
            "knowledge_base_ids": [base["id"]],
            "generate_answer": True,
        },
    )
    assert response.status_code == 200
    events = [
        json.loads(line.removeprefix("data: "))
        for line in response.text.splitlines()
        if line.startswith("data: ")
    ]
    step_types = [event["step"]["type"] for event in events if event["type"] == "step"]
    expected = [
        "stream_connected",
        "scope_resolved",
        "query_rewrite_started",
        "query_rewritten",
        "hybrid_retrieval_started",
        "hybrid_retrieval_completed",
        "fusion_completed",
        "rerank_started",
        "rerank_completed",
        "context_assembled",
        "answer_generation_started",
        "answer_generated",
    ]
    assert [name for name in step_types if name != "knowledge_waiting"] == expected
    result_event = next(event for event in events if event["type"] == "knowledge_result")
    assert result_event["result"]["chunks"]
    assert events[-1]["type"] == "done"


def test_duplicate_document_is_not_counted_twice(client):
    base = _new_base(client)
    payload = {"title": "重复资料", "content": "完全相同的有效知识内容。" * 20, "source": "测试"}
    first = client.post(f"/api/knowledge-bases/{base['id']}/documents/text", json=payload)
    second = client.post(f"/api/knowledge-bases/{base['id']}/documents/text", json=payload)
    assert first.status_code == second.status_code == 201
    assert first.json()["id"] == second.json()["id"]
    bases = client.get("/api/knowledge-bases").json()
    current = next(item for item in bases if item["id"] == base["id"])
    assert current["document_count"] == 1


def test_folder_upload_preserves_relative_path_and_cleans_duplicate_source(client):
    base = _new_base(client)
    content = ("文件夹批量导入需要保留相对目录，并对重复资料进行内容去重。\n" * 20).encode()
    upload = {
        "files": {"file": ("guide.md", content, "text/markdown")},
        "data": {"relative_path": "课程资料/第一章/guide.md"},
    }

    first = client.post(
        f"/api/knowledge-bases/{base['id']}/documents/upload", **upload
    )
    duplicate = client.post(
        f"/api/knowledge-bases/{base['id']}/documents/upload", **upload
    )

    assert first.status_code == duplicate.status_code == 201
    assert first.json()["title"] == "课程资料/第一章/guide.md"
    assert first.json()["source"] == "本地文件：课程资料/第一章/guide.md"
    assert duplicate.json()["id"] == first.json()["id"]
    assert duplicate.json()["ingestion"]["duplicate"] is True
    detail = client.get(f"/api/knowledge-documents/{first.json()['id']}").json()
    assert detail["metadata"]["relative_path"] == "课程资料/第一章/guide.md"
    sources = client.get(f"/api/knowledge-bases/{base['id']}/sources").json()
    assert len(sources) == 1
    assert sources[0]["uri"] == "课程资料/第一章/guide.md"


def test_database_source_full_sync(client, tmp_path):
    path = tmp_path / "source.db"
    connection = sqlite3.connect(path)
    connection.execute("CREATE TABLE papers(title TEXT, abstract TEXT)")
    connection.execute(
        "INSERT INTO papers VALUES (?, ?)",
        ("Mesh quality", "Jacobian determinant and orthogonality are core metrics."),
    )
    connection.commit()
    connection.close()
    base = _new_base(client)
    response = client.post(
        f"/api/knowledge-bases/{base['id']}/sources/database",
        json={
            "name": "论文数据库",
            "connection_url": f"sqlite:///{path.as_posix()}",
            "query": "SELECT title, abstract FROM papers",
            "sync_now": True,
        },
    )
    assert response.status_code == 201, response.text
    body = response.json()
    assert body["job"]["status"] == "completed"
    assert body["job"]["document_count"] == 1
    sources = client.get(f"/api/knowledge-bases/{base['id']}/sources").json()
    assert sources[0]["status"] == "ready"
    assert "config_ciphertext" not in sources[0]


def test_external_sources_can_be_registered_without_immediate_network(client):
    base = _new_base(client)
    web = client.post(
        f"/api/knowledge-bases/{base['id']}/sources/web",
        json={"name": "规范站点", "url": "https://example.com/standards", "sync_now": False},
    )
    api = client.post(
        f"/api/knowledge-bases/{base['id']}/sources/api",
        json={
            "name": "第三方数据",
            "url": "https://example.com/api/data",
            "headers": {"Authorization": "Bearer secret"},
            "response_path": "data.items",
            "sync_now": False,
        },
    )
    assert web.status_code == api.status_code == 201
    listed = client.get(f"/api/knowledge-bases/{base['id']}/sources").json()
    assert {item["source_type"] for item in listed} == {"web", "api"}
    assert all("config_ciphertext" not in item for item in listed)


def test_knowledge_group_crud_and_scoped_retrieval(client):
    aerodynamics = _new_base(client)
    medicine = _new_base(client)
    for base, title, content in (
        (aerodynamics, "空气动力学", "机翼升力与压力分布是空气动力学分析的重要内容。" * 12),
        (medicine, "细胞医学", "细胞膜蛋白与免疫反应是医学研究的重要内容。" * 12),
    ):
        response = client.post(
            f"/api/knowledge-bases/{base['id']}/documents/text",
            json={"title": title, "content": content, "source": "分组测试"},
        )
        assert response.status_code == 201

    created = client.post(
        "/api/knowledge-groups",
        json={
            "name": f"航空学科-{uuid.uuid4()}",
            "description": "仅包含航空知识库",
            "color": "#2878c8",
            "knowledge_base_ids": [aerodynamics["id"]],
        },
    )
    assert created.status_code == 201
    group = created.json()
    assert group["knowledge_base_ids"] == [aerodynamics["id"]]

    scoped = client.post(
        "/api/knowledge/query",
        json={
            "query": "细胞免疫",
            "knowledge_group_ids": [group["id"]],
            "generate_answer": False,
        },
    )
    assert scoped.status_code == 200
    assert scoped.json()["trace"]["knowledge_base_ids"] == [aerodynamics["id"]]
    assert all(item["knowledge_base_id"] == aerodynamics["id"] for item in scoped.json()["chunks"])

    updated = client.put(
        f"/api/knowledge-groups/{group['id']}/members",
        json={"knowledge_base_ids": [medicine["id"]]},
    )
    assert updated.status_code == 200
    assert updated.json()["knowledge_base_ids"] == [medicine["id"]]

    rescoped = client.post(
        "/api/knowledge/query",
        json={
            "query": "细胞免疫",
            "knowledge_group_ids": [group["id"]],
            "generate_answer": False,
        },
    ).json()
    assert rescoped["chunks"]
    assert all(item["knowledge_base_id"] == medicine["id"] for item in rescoped["chunks"])

    deleted = client.delete(f"/api/knowledge-groups/{group['id']}")
    assert deleted.status_code == 204
    assert client.get(f"/api/knowledge-bases/{medicine['id']}/documents").status_code == 200


def test_knowledge_internal_crud_rebuilds_indexes(client):
    base = _new_base(client)
    renamed = f"可编辑知识库-{uuid.uuid4()}"
    base_update = client.patch(
        f"/api/knowledge-bases/{base['id']}",
        json={"name": renamed, "discipline": "材料学", "description": "CRUD 测试"},
    )
    assert base_update.status_code == 200
    assert base_update.json()["name"] == renamed

    created = client.post(
        f"/api/knowledge-bases/{base['id']}/documents/text",
        json={"title": "旧标题", "source": "旧来源", "content": "旧正文描述金属材料。" * 20},
    )
    assert created.status_code == 201
    old_document_id = created.json()["id"]

    metadata_update = client.patch(
        f"/api/knowledge-documents/{old_document_id}",
        json={"title": "材料规范", "source": "实验室规范"},
    )
    assert metadata_update.status_code == 200
    assert metadata_update.json()["title"] == "材料规范"

    content_update = client.patch(
        f"/api/knowledge-documents/{old_document_id}",
        json={
            "title": "材料规范新版",
            "source": "实验室规范 V2",
            "content": "钛合金疲劳裂纹扩展速率需要通过循环载荷试验评价。" * 20,
        },
    )
    assert content_update.status_code == 200, content_update.text
    new_document_id = content_update.json()["id"]
    assert new_document_id != old_document_id
    assert client.get(f"/api/knowledge-documents/{old_document_id}").status_code == 404
    detail = client.get(f"/api/knowledge-documents/{new_document_id}").json()
    assert "钛合金疲劳裂纹" in detail["cleaned_content"]
    assert detail["embedding_count"] >= 1

    source_created = client.post(
        f"/api/knowledge-bases/{base['id']}/sources/web",
        json={"name": "待编辑网页", "url": "https://example.com/source", "sync_now": False},
    ).json()["source"]
    source_updated = client.patch(
        f"/api/knowledge-sources/{source_created['id']}",
        json={"name": "已编辑网页", "uri": "https://example.com/new-source?token=hidden"},
    )
    assert source_updated.status_code == 200
    assert source_updated.json()["name"] == "已编辑网页"
    assert "token=" not in source_updated.json()["uri"]
    assert client.delete(f"/api/knowledge-sources/{source_created['id']}").status_code == 204

    assert client.delete(f"/api/knowledge-documents/{new_document_id}").status_code == 204
    overview = client.get(f"/api/knowledge-bases/{base['id']}/overview").json()
    assert overview["statistics"]["documents"] == 0
    assert overview["statistics"]["embeddings"] == 0
    assert client.delete(f"/api/knowledge-bases/{base['id']}").status_code == 204
    assert client.get(f"/api/knowledge-bases/{base['id']}/overview").status_code == 404


@pytest.mark.asyncio
async def test_siliconflow_embedding_and_rerank_payloads(monkeypatch):
    calls = []

    class FakeResponse:
        is_error = False
        status_code = 200
        text = ""

        def raise_for_status(self):
            return None

        def json(self):
            if calls[-1][0].endswith("embeddings"):
                return {"data": [{"index": 0, "embedding": [0.1, 0.2, 0.3]}]}
            return {"results": [{"index": 1, "relevance_score": 0.98}]}

    class FakeClient:
        def __init__(self, *args, **kwargs):
            pass

        async def __aenter__(self):
            return self

        async def __aexit__(self, *args):
            return None

        async def post(self, url, **kwargs):
            calls.append((url, kwargs))
            return FakeResponse()

    monkeypatch.setattr("backend.app.services.knowledge_vector.httpx.AsyncClient", FakeClient)
    config = KnowledgeProviderConfig(
        id="mock",
        api_key_ciphertext=secret_store.encrypt("test-only-key"),
        embedding_base_url="https://mock.local/v1/embeddings",
        rerank_base_url="https://mock.local/v1/rerank",
    )
    vectors = await EmbeddingClient(config).embed(["网格质量"])
    ranks = await RerankClient(config).rerank("质量", ["无关", "网格质量"], 1)
    assert vectors == [[0.1, 0.2, 0.3]]
    assert ranks == [(1, 0.98)]
    assert calls[0][1]["json"]["model"] == "Qwen/Qwen3-VL-Embedding-8B"
    assert calls[1][1]["json"]["model"] == "BAAI/bge-reranker-v2-m3"
    assert calls[1][1]["json"]["top_n"] == 1
