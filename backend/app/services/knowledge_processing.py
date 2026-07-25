from __future__ import annotations

import csv
import hashlib
import io
import re
import unicodedata
from collections import Counter
from dataclasses import dataclass, field
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


@dataclass
class ExtractedSection:
    text: str
    heading: str = ""
    locator: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass
class ChunkDraft:
    content: str
    level: str
    parent_index: int | None
    metadata: dict[str, Any]


class MainTextHTMLParser(HTMLParser):
    """Small dependency-free HTML extractor that removes navigation and executable content."""

    ignored = {"script", "style", "noscript", "svg", "nav", "footer", "header", "form"}
    blocks = {"p", "div", "article", "section", "main", "li", "br", "tr", "h1", "h2", "h3"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.depth = 0
        self.parts: list[str] = []
        self.title = ""
        self._in_title = False
        self.links: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        tag = tag.lower()
        if tag == "a":
            href = dict(attrs).get("href")
            if href:
                self.links.append(href)
        if tag in self.ignored:
            self.depth += 1
        if tag == "title":
            self._in_title = True
        if not self.depth and tag in self.blocks:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag == "title":
            self._in_title = False
        if tag in self.ignored and self.depth:
            self.depth -= 1
        if not self.depth and tag in self.blocks:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if self.depth:
            return
        if self._in_title:
            self.title += data
        self.parts.append(data)

    @property
    def text(self) -> str:
        return "".join(self.parts)


def html_to_text(value: str) -> tuple[str, str]:
    parser = MainTextHTMLParser()
    parser.feed(value)
    return parser.text, clean_text(parser.title)[0]


def clean_text(value: str) -> tuple[str, dict[str, int]]:
    """Normalize Unicode, remove control/noise lines, and collapse repeated lines."""

    original_length = len(value)
    value = unicodedata.normalize("NFKC", value.replace("\u00a0", " "))
    value = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", value)
    raw_lines = [re.sub(r"[ \t]+", " ", line).strip() for line in value.splitlines()]
    counts = Counter(line for line in raw_lines if 0 < len(line) <= 120)
    seen: Counter[str] = Counter()
    result: list[str] = []
    removed_repeated = 0
    removed_noise = 0
    for line in raw_lines:
        if not line:
            if result and result[-1]:
                result.append("")
            continue
        compact = re.sub(r"\s+", "", line)
        if re.fullmatch(r"(?:[-_=*•·.。]{3,}|\d{1,4})", compact):
            removed_noise += 1
            continue
        # Repeated short headers/footers are retained once, then treated as layout noise.
        if counts[line] >= 3:
            seen[line] += 1
            if seen[line] > 1:
                removed_repeated += 1
                continue
        result.append(line)
    cleaned = "\n".join(result)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned).strip()
    return cleaned, {
        "original_chars": original_length,
        "cleaned_chars": len(cleaned),
        "noise_lines_removed": removed_noise,
        "repeated_lines_removed": removed_repeated,
    }


def extract_sections(filename: str, data: bytes) -> tuple[list[ExtractedSection], str]:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        reader = PdfReader(io.BytesIO(data))
        sections = [
            ExtractedSection(
                text=page.extract_text() or "",
                locator=f"第 {index} 页",
                metadata={"page": index},
            )
            for index, page in enumerate(reader.pages, 1)
        ]
        return sections, "application/pdf"
    if suffix == ".docx":
        document = Document(io.BytesIO(data))
        sections: list[ExtractedSection] = []
        heading = ""
        buffer: list[str] = []
        section_index = 1
        list_counter = 0
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = (paragraph.style.name if paragraph.style else "").lower()
            if "heading" in style_name or "标题" in style_name:
                if buffer:
                    sections.append(
                        ExtractedSection("\n".join(buffer), heading, f"章节 {section_index}")
                    )
                    section_index += 1
                    buffer = []
                heading = text
                list_counter = 0
            else:
                paragraph_properties = getattr(paragraph._p, "pPr", None)
                numbering = getattr(paragraph_properties, "numPr", None)
                is_list = numbering is not None or "list" in style_name or "列表" in style_name
                if is_list:
                    list_counter += 1
                    text = f"{list_counter}. {text}"
                else:
                    list_counter = 0
                buffer.append(text)
        if buffer or heading:
            sections.append(ExtractedSection("\n".join(buffer), heading, f"章节 {section_index}"))
        for table_index, table in enumerate(document.tables, 1):
            rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
            sections.append(
                ExtractedSection("\n".join(rows), f"表格 {table_index}", f"表格 {table_index}")
            )
        return sections, "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if suffix == ".pptx":
        presentation = Presentation(io.BytesIO(data))
        sections = []
        for slide_index, slide in enumerate(presentation.slides, 1):
            parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and str(shape.text).strip():
                    parts.append(str(shape.text).strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        parts.append(" | ".join(cell.text.strip() for cell in row.cells))
            title = ""
            if slide.shapes.title is not None:
                title = slide.shapes.title.text.strip()
            sections.append(
                ExtractedSection("\n".join(parts), title, f"第 {slide_index} 页幻灯片", {"slide": slide_index})
            )
        return sections, "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if suffix in {".txt", ".md", ".csv", ".json", ".html", ".htm"}:
        decoded = data.decode("utf-8-sig", errors="replace")
        if suffix in {".html", ".htm"}:
            decoded, title = html_to_text(decoded)
            return [ExtractedSection(decoded, title)], "text/html"
        if suffix == ".csv":
            rows = list(csv.reader(io.StringIO(decoded)))
            decoded = "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)
        return [ExtractedSection(decoded)], "text/plain"
    if suffix == ".ppt":
        raise ValueError("旧版 .ppt 需先另存为 .pptx 后上传")
    if suffix == ".doc":
        raise ValueError("旧版 .doc 需先另存为 .docx 后上传")
    raise ValueError("支持 PDF、DOCX、PPTX、TXT、MD、CSV、JSON 和 HTML 文件")


def _sentences(text: str) -> list[str]:
    pieces = re.split(r"(?<=[。！？!?；;\.])\s+|\n+", text)
    return [piece.strip() for piece in pieces if piece.strip()]


LIST_ITEM_PATTERN = re.compile(
    r"^\s*(?P<marker>(?:\d{1,3}|[一二三四五六七八九十百]+)[、.．)]|"
    r"（(?:\d{1,3}|[一二三四五六七八九十百]+)）|[①②③④⑤⑥⑦⑧⑨⑩]|[-•·])\s*(?P<body>.+)$"
)


@dataclass
class StructuredBlock:
    kind: str
    text: str = ""
    intro: str = ""
    items: list[tuple[str, str]] = field(default_factory=list)


@dataclass
class StructuredEntry:
    block: StructuredBlock
    start_section_index: int
    end_section_index: int


@dataclass
class NumberedListSpan:
    content: str
    intro: str
    items: list[tuple[str, str]]
    start_section_index: int
    end_section_index: int
    locators: list[str]


_CIRCLED_ORDINALS = {value: index for index, value in enumerate("①②③④⑤⑥⑦⑧⑨⑩", 1)}
_CHINESE_DIGITS = {"零": 0, "一": 1, "二": 2, "三": 3, "四": 4, "五": 5, "六": 6, "七": 7, "八": 8, "九": 9}


def _marker_ordinal(marker: str) -> int | None:
    """Return a comparable ordinal for common list markers; bullets have no ordinal."""

    marker = marker.strip()
    if marker in _CIRCLED_ORDINALS:
        return _CIRCLED_ORDINALS[marker]
    digits = re.search(r"\d{1,3}", marker)
    if digits:
        return int(digits.group())
    value = re.sub(r"[、.．()（）\s]", "", marker)
    if not value or any(char not in _CHINESE_DIGITS and char not in {"十", "百"} for char in value):
        return None
    if value == "十":
        return 10
    if "百" in value:
        head, _, tail = value.partition("百")
        hundreds = _CHINESE_DIGITS.get(head, 1) if head else 1
        remainder = _marker_ordinal(tail) if tail else 0
        return hundreds * 100 + (remainder or 0)
    if "十" in value:
        head, _, tail = value.partition("十")
        tens = _CHINESE_DIGITS.get(head, 1) if head else 1
        ones = _CHINESE_DIGITS.get(tail, 0) if tail else 0
        return tens * 10 + ones
    return _CHINESE_DIGITS.get(value)


def _nested_intro_start(lines: list[str]) -> int | None:
    """Locate a likely sub-section heading that was wrapped into the prior list item."""

    for index, line in enumerate(lines[1:], 1):
        value = line.strip()
        if len(value) > 80:
            continue
        if "★" in value or re.search(
            r"(?:原理|流程|机制|方法|管理|映射|算法|特点|优缺点|基本概念|经典问题)"
            r"(?:[-—:：、（(]|$)",
            value,
        ):
            return index
    return None


def _structured_blocks(text: str) -> list[StructuredBlock]:
    """Preserve contiguous numbered lists instead of packing every line as plain prose."""

    blocks: list[StructuredBlock] = []
    normal_lines: list[str] = []
    list_items: list[tuple[str, str]] = []
    current_marker = ""
    current_lines: list[str] = []
    list_intro = ""
    list_gap = False

    def flush_normal() -> None:
        value = "\n".join(normal_lines).strip()
        if value:
            blocks.append(StructuredBlock("text", text=value))
        normal_lines.clear()

    def flush_item() -> None:
        nonlocal current_marker, current_lines, list_gap
        value = "\n".join(current_lines).strip()
        if current_marker and value:
            list_items.append((current_marker, value))
        current_marker = ""
        current_lines = []
        list_gap = False

    def flush_list() -> None:
        nonlocal list_intro, list_items
        flush_item()
        if len(list_items) >= 2:
            blocks.append(StructuredBlock("numbered_list", intro=list_intro, items=list_items))
        elif list_items:
            marker, value = list_items[0]
            blocks.append(StructuredBlock("text", text=f"{marker} {value}"))
        list_intro = ""
        list_items = []

    for line in text.splitlines():
        stripped = line.strip()
        match = LIST_ITEM_PATTERN.match(stripped)
        if match:
            next_marker = match.group("marker")
            current_ordinal = _marker_ordinal(current_marker) if current_marker else None
            next_ordinal = _marker_ordinal(next_marker)
            if (
                current_marker
                and current_ordinal is not None
                and next_ordinal is not None
                and next_ordinal <= current_ordinal
            ):
                # A reset such as 1..4 followed by a nested/new 1..5 list must not
                # become one oversized list. Preserve the trailing outer item as
                # the new list's local introduction when it contains sub-section text.
                nested_intro = ""
                if len(current_lines) >= 2:
                    intro_start = _nested_intro_start(current_lines)
                    if intro_start is not None:
                        nested_intro = "\n".join(current_lines[intro_start:]).strip()
                        current_lines = current_lines[:intro_start]
                    else:
                        nested_intro = "\n".join(
                            [f"{current_marker} {current_lines[0]}", *current_lines[1:]]
                        ).strip()
                flush_list()
                list_intro = nested_intro
            if not current_marker and not list_items:
                if not list_intro:
                    list_intro = next((item for item in reversed(normal_lines) if item.strip()), "")
                flush_normal()
            flush_item()
            current_marker = next_marker
            current_lines = [match.group("body").strip()]
            continue
        if current_marker:
            if not stripped:
                list_gap = True
            elif list_gap:
                flush_list()
                normal_lines.append(line)
            else:
                current_lines.append(stripped)
            continue
        normal_lines.append(line)

    if current_marker or list_items:
        flush_list()
    flush_normal()
    return blocks


def _structured_entries(sections: list[ExtractedSection]) -> list[StructuredEntry]:
    """Parse all sections and stitch a numbered list continued on the next page/section."""

    entries: list[StructuredEntry] = []
    for section_index, section in enumerate(sections):
        cleaned, _stats = clean_text(section.text)
        if not cleaned:
            continue
        blocks = _structured_blocks(cleaned)
        if (
            entries
            and blocks
            and entries[-1].end_section_index == section_index - 1
            and entries[-1].block.kind == "numbered_list"
            and blocks[0].kind == "numbered_list"
        ):
            previous_items = entries[-1].block.items
            current_items = blocks[0].items
            previous_ordinal = _marker_ordinal(previous_items[-1][0]) if previous_items else None
            current_ordinal = _marker_ordinal(current_items[0][0]) if current_items else None
            if (
                previous_ordinal is not None
                and current_ordinal is not None
                and current_ordinal == previous_ordinal + 1
            ):
                entries[-1].block.items.extend(current_items)
                entries[-1].end_section_index = section_index
                blocks = blocks[1:]
        entries.extend(
            StructuredEntry(block, section_index, section_index)
            for block in blocks
        )
    return entries


def numbered_list_spans(sections: list[ExtractedSection]) -> list[NumberedListSpan]:
    """Expose focused, cross-section list contexts for exhaustive retrieval."""

    spans: list[NumberedListSpan] = []
    for entry in _structured_entries(sections):
        if entry.block.kind != "numbered_list":
            continue
        rendered_items = "\n".join(f"{marker} {value}" for marker, value in entry.block.items)
        content = "\n".join(
            value for value in (entry.block.intro.strip(), rendered_items) if value
        )
        locators = [
            sections[index].locator
            for index in range(entry.start_section_index, entry.end_section_index + 1)
            if sections[index].locator
        ]
        spans.append(
            NumberedListSpan(
                content=content,
                intro=entry.block.intro.strip(),
                items=list(entry.block.items),
                start_section_index=entry.start_section_index,
                end_section_index=entry.end_section_index,
                locators=list(dict.fromkeys(locators)),
            )
        )
    return spans


def _pack_units(units: list[str], target: int, overlap: int) -> list[str]:
    chunks: list[str] = []
    buffer: list[str] = []
    length = 0
    for unit in units:
        if len(unit) > target:
            if buffer:
                chunks.append("\n".join(buffer))
                buffer, length = [], 0
            step = max(1, target - overlap)
            chunks.extend(unit[start : start + target] for start in range(0, len(unit), step))
            continue
        if buffer and length + len(unit) + 1 > target:
            packed = "\n".join(buffer)
            chunks.append(packed)
            tail = packed[-overlap:] if overlap else ""
            buffer = [tail, unit] if tail else [unit]
            length = sum(map(len, buffer)) + len(buffer) - 1
        else:
            buffer.append(unit)
            length += len(unit) + 1
    if buffer:
        chunks.append("\n".join(buffer))
    return [chunk.strip() for chunk in chunks if chunk.strip()]


def hierarchical_chunks(
    sections: list[ExtractedSection],
    *,
    parent_size: int = 1600,
    child_size: int = 480,
    child_overlap: int = 80,
) -> tuple[list[ChunkDraft], dict[str, int]]:
    """Structure-aware parent/child segmentation with sentence-boundary overlap."""

    drafts: list[ChunkDraft] = []
    fingerprints: set[str] = set()
    duplicate_chunks = 0
    numbered_lists = 0
    numbered_list_items = 0
    parent_index = 0
    list_numbers: Counter[int] = Counter()
    for entry in _structured_entries(sections):
        section_index = entry.start_section_index
        section = sections[section_index]
        block = entry.block
        prefix = f"# {section.heading}\n\n" if section.heading else ""
        base_meta = {**section.metadata, "heading": section.heading, "locator": section.locator}
        if entry.end_section_index > section_index:
            end_section = sections[entry.end_section_index]
            locators = [
                sections[index].locator
                for index in range(section_index, entry.end_section_index + 1)
                if sections[index].locator
            ]
            unique_locators = list(dict.fromkeys(locators))
            base_meta.update(
                {
                    "locator": " – ".join(unique_locators),
                    "locators": unique_locators,
                    "cross_section_continuation": True,
                    "section_start": section_index + 1,
                    "section_end": entry.end_section_index + 1,
                }
            )
            if "page" in section.metadata and "page" in end_section.metadata:
                base_meta.update(
                    {
                        "page_start": section.metadata["page"],
                        "page_end": end_section.metadata["page"],
                    }
                )
        if block.kind == "numbered_list":
            list_numbers[section_index] += 1
            numbered_lists += 1
            numbered_list_items += len(block.items)
            list_id = f"section-{section_index + 1}-list-{list_numbers[section_index]}"
            rendered_items = "\n".join(f"{marker} {value}" for marker, value in block.items)
            parent_body = "\n".join(
                value for value in (block.intro.strip(), rendered_items) if value
            )
            parent_content = f"{prefix}{parent_body}".strip()
            parent_meta = {
                **base_meta,
                "structure": "numbered_list",
                "list_id": list_id,
                "list_item_count": len(block.items),
            }
            drafts.append(ChunkDraft(parent_content, "parent", None, parent_meta))
            intro = block.intro.strip()
            for item_index, (marker, item_text) in enumerate(block.items, 1):
                shared_prefix = "\n".join(value for value in (prefix.strip(), intro) if value)
                available = max(160, child_size - len(shared_prefix) - len(marker) - 2)
                item_parts = _pack_units(_sentences(item_text), available, child_overlap)
                for part_index, part in enumerate(item_parts, 1):
                    child_content = "\n".join(
                        value for value in (shared_prefix, f"{marker} {part}") if value
                    ).strip()
                    fingerprint = hashlib.sha256(
                        re.sub(r"[\W_]", "", child_content).lower().encode("utf-8")
                    ).hexdigest()
                    if fingerprint in fingerprints:
                        duplicate_chunks += 1
                        continue
                    fingerprints.add(fingerprint)
                    drafts.append(
                        ChunkDraft(
                            child_content,
                            "child",
                            parent_index,
                            {
                                **parent_meta,
                                "list_item_index": item_index,
                                "list_item_marker": marker,
                                "list_item_part": part_index,
                            },
                        )
                    )
            parent_index += 1
            continue

        parents = _pack_units(_sentences(block.text), parent_size, 0)
        for parent in parents:
            parent_content = f"{prefix}{parent}".strip()
            drafts.append(ChunkDraft(parent_content, "parent", None, base_meta))
            children = _pack_units(_sentences(parent), child_size, child_overlap)
            for child in children:
                child_content = f"{prefix}{child}".strip()
                fingerprint = hashlib.sha256(
                    re.sub(r"[\W_]", "", child_content).lower().encode("utf-8")
                ).hexdigest()
                if fingerprint in fingerprints:
                    duplicate_chunks += 1
                    continue
                fingerprints.add(fingerprint)
                drafts.append(ChunkDraft(child_content, "child", parent_index, base_meta))
            parent_index += 1
    return drafts, {
        "duplicate_chunks_removed": duplicate_chunks,
        "numbered_lists": numbered_lists,
        "numbered_list_items": numbered_list_items,
    }


def estimate_tokens(text: str) -> int:
    chinese = len(re.findall(r"[\u3400-\u9fff]", text))
    other = max(0, len(text) - chinese)
    return max(1, chinese + other // 4)
